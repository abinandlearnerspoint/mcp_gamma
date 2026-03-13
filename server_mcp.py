#!/usr/bin/env python3
"""
Gamma Presentation MCP Server
Supports: generate, from-template, themes, folders, status check
"""

import os
import sys
from typing import Any, Dict, List, Literal, Optional, Tuple
import asyncio
import httpx
from dotenv import load_dotenv
from fastmcp import FastMCP
from pptx import Presentation

load_dotenv()

# ── Config ─────────────────────────────────────────────────────────────────
GAMMA_API_BASE    = "https://public-api.gamma.app/v1.0"
GAMMA_GENERATIONS = f"{GAMMA_API_BASE}/generations"
GAMMA_API_KEY     = os.getenv("GAMMA_API_KEY")

# Polling — Gamma can take 2–5 minutes for complex decks
POLL_INTERVAL_SEC   = 8    # seconds between polls
POLL_MAX_ATTEMPTS   = 45   # 45 × 8s = 6 minutes max
POLL_MAX_ERR_STREAK = 3    # consecutive HTTP errors before giving up

# ── FastMCP ─────────────────────────────────────────────────────────────────
mcp = FastMCP("gamma-presentation")


# ═══════════════════════════════════════════════════════════════════════════
# INTERNAL HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def _headers() -> Dict[str, str]:
    return {
        "Content-Type": "application/json",
        "X-API-KEY": GAMMA_API_KEY or "",
    }


def _check_api_key() -> Optional[str]:
    """Return an error string if API key is missing, else None."""
    if not GAMMA_API_KEY:
        return "❌ GAMMA_API_KEY is not set. Add it to your .env file."
    return None


def _extract_gamma_error(response: httpx.Response) -> str:
    """Best-effort extraction of a readable error message from a Gamma error response."""
    try:
        body = response.json()
        return body.get("message") or body.get("error") or response.text
    except Exception:
        return response.text or f"HTTP {response.status_code}"


def _extract_url(data: Dict[str, Any]) -> Optional[str]:
    """
    Try all known URL field names in a completed generation response.
    Returns None if no URL found — callers must handle this explicitly.
    No fallback URL is constructed: a guessed URL is worse than no URL.
    """
    return (
        data.get("url")
        or data.get("link")
        or data.get("gammaUrl")
        or (data.get("gamma") or {}).get("url")
        or (data.get("result") or {}).get("url")
        or (data.get("data") or {}).get("url")
        or (data.get("output") or {}).get("url")
    )


def _format_result(result: Dict[str, Optional[str]], label: str) -> str:
    """Format the final user-facing response for a completed generation."""
    url      = result.get("url")
    pdf_url  = result.get("pdfUrl")
    pptx_url = result.get("pptxUrl")
    error    = result.get("error")

    if not url:
        return (
            f"❌ {label} failed.\n"
            f"Error: {error or 'Unknown error'}\n\n"
            f"💡 The generation may still complete — "
            f"check https://gamma.app dashboard."
        )

    lines = [f"✅ {label} generated successfully!", "", f"🔗 View here : {url}"]
    if pdf_url:
        lines.append(f"📄 PDF export : {pdf_url}  ⚠️  Download soon — link expires!")
    if pptx_url:
        lines.append(f"📊 PPTX export: {pptx_url}  ⚠️  Download soon — link expires!")
    return "\n".join(lines)


# ── Core: POST to Gamma, return generation_id ───────────────────────────────
# Returns a 3-tuple:
#   ("direct", url,          None)   → URL returned immediately, skip polling
#   ("poll",   generation_id, None)  → Need to poll
#   ("error",  None,          msg)   → Request failed
async def _start_generation(
    endpoint: str,
    params: Dict[str, Any],
    label: str = "generation",
) -> Tuple[str, Optional[str], Optional[str]]:

    print(f"\n{'='*60}", flush=True)
    print(f"📤 Starting {label}", flush=True)
    print(f"🔗 Endpoint : {endpoint}", flush=True)
    print(f"📦 Payload  : {params}", flush=True)
    print(f"{'='*60}", flush=True)

    try:
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(endpoint, headers=_headers(), json=params)

        print(f"📥 HTTP {resp.status_code} | {resp.text[:800]}", flush=True)

        if not resp.is_success:
            return "error", None, f"HTTP {resp.status_code}: {_extract_gamma_error(resp)}"

        data = resp.json()
        print(f"📥 Full Response: {data}", flush=True)

        generation_id = data.get("generationId") or data.get("id")

        if not generation_id:
            direct_url = data.get("url") or data.get("link")
            if direct_url:
                return "direct", direct_url, None
            return "error", None, f"No generationId in response: {data}"

        print(f"🆔 Generation ID: {generation_id}", flush=True)
        return "poll", generation_id, None

    except httpx.TimeoutException:
        return "error", None, "Request to Gamma API timed out. Check your network."
    except httpx.ConnectError:
        return "error", None, "Cannot connect to Gamma API. Check your network."
    except Exception as e:
        return "error", None, f"{type(e).__name__}: {e}"


# ── Core: Poll until done ────────────────────────────────────────────────────
async def _poll_generation(generation_id: str) -> Dict[str, Optional[str]]:
    """
    Poll until status = completed / failed, or timeout.
    Returns dict with keys: url, pdfUrl, pptxUrl, error
    """
    status_url = f"{GAMMA_GENERATIONS}/{generation_id}"
    total_wait = POLL_MAX_ATTEMPTS * POLL_INTERVAL_SEC
    err_streak = 0

    print(f"\n⏳ Polling: {status_url}", flush=True)
    print(f"⏳ Max wait: {total_wait}s ({POLL_MAX_ATTEMPTS} × {POLL_INTERVAL_SEC}s)", flush=True)

    def _fail(msg: str) -> Dict[str, Optional[str]]:
        return {"url": None, "pdfUrl": None, "pptxUrl": None, "error": msg}

    for attempt in range(1, POLL_MAX_ATTEMPTS + 1):
        await asyncio.sleep(POLL_INTERVAL_SEC)
        elapsed = attempt * POLL_INTERVAL_SEC

        try:
            # Fresh client per poll — avoids stale connections on long waits
            async with httpx.AsyncClient(timeout=30) as client:
                r = await client.get(status_url, headers=_headers())

            print(
                f"🔄 Poll {attempt:02d}/{POLL_MAX_ATTEMPTS} "
                f"| {elapsed}s elapsed | HTTP {r.status_code} | {r.text[:300]}",
                flush=True,
            )

            if not r.is_success:
                err_streak += 1
                print(f"⚠️  Error streak: {err_streak}/{POLL_MAX_ERR_STREAK}", flush=True)
                if err_streak >= POLL_MAX_ERR_STREAK:
                    return _fail(
                        f"{POLL_MAX_ERR_STREAK} consecutive poll failures. "
                        f"Last: HTTP {r.status_code}: {r.text[:200]}"
                    )
                continue

            err_streak = 0
            data   = r.json()
            status = data.get("status", "").lower().strip()
            print(f"📊 Status: '{status}'", flush=True)

            # ── Completed ──────────────────────────────────────────────────
            if status in ("completed", "complete", "done", "success"):
                url = _extract_url(data)
                if not url:
                    print(f"⚠️  Completed but no URL. Full data: {data}", flush=True)
                    return _fail(
                        f"Generation completed but Gamma returned no URL. "
                        f"Check https://gamma.app — your deck may be there. "
                        f"Response keys: {list(data.keys())}"
                    )

                pdf_url  = data.get("pdfUrl")  or (data.get("exports") or {}).get("pdf")
                pptx_url = data.get("pptxUrl") or (data.get("exports") or {}).get("pptx")

                print(f"✅ Done! URL={url} PDF={pdf_url} PPTX={pptx_url}", flush=True)
                return {"url": url, "pdfUrl": pdf_url, "pptxUrl": pptx_url, "error": None}

            # ── Failed ─────────────────────────────────────────────────────
            elif status in ("failed", "error", "cancelled", "canceled"):
                msg = data.get("message") or data.get("error") or data.get("reason") or status
                print(f"❌ Generation failed: {msg}", flush=True)
                return _fail(f"Generation failed (status='{status}'): {msg}")

            # ── Still running ──────────────────────────────────────────────
            else:
                remaining = total_wait - elapsed
                print(f"⏳ Still {status or 'pending'}... ~{remaining}s remaining", flush=True)

        except httpx.TimeoutException:
            err_streak += 1
            print(f"⚠️  Poll timeout (streak {err_streak})", flush=True)
            if err_streak >= POLL_MAX_ERR_STREAK:
                return _fail("Repeated poll timeouts. Check your network.")

        except Exception as e:
            err_streak += 1
            print(f"⚠️  Poll exception: {e} (streak {err_streak})", flush=True)
            if err_streak >= POLL_MAX_ERR_STREAK:
                return _fail(f"Repeated poll errors: {type(e).__name__}: {e}")

    # ── Timeout ─────────────────────────────────────────────────────────────
    return _fail(
        f"Timed out after {total_wait}s ({POLL_MAX_ATTEMPTS} polls). "
        f"Generation may still complete — check https://gamma.app dashboard. "
        f"Or use check_generation_status('{generation_id}') to resume."
    )


# ═══════════════════════════════════════════════════════════════════════════
# TOOL 1 — Generate Presentation (new, from scratch)
# ═══════════════════════════════════════════════════════════════════════════
@mcp.tool()
async def generate_presentation(
    inputText: str,
    numCards: Optional[int] = 5,
    textMode: Optional[Literal["generate", "summarize"]] = "generate",
    additionalInstructions: Optional[str] = None,
) -> str:
    """
    Generate a brand-new Gamma presentation from a topic or prompt.

    Args:
        inputText              : Topic or content for the presentation (REQUIRED)
        numCards               : Number of slides 1–20 (default: 5)
        textMode               : 'generate' to create new content,
                                 'summarize' to condense your input text
                                 (default: 'generate')
        additionalInstructions : Extra instructions e.g. 'Use a formal tone'

    Returns:
        Link to the generated presentation.
    """
    print(f"\n🚀 generate_presentation()", flush=True)
    print(f"   inputText : {inputText[:80]}", flush=True)
    print(f"   numCards  : {numCards}", flush=True)

    err = _check_api_key()
    if err:
        return err

    if not inputText or not inputText.strip():
        return "❌ inputText cannot be empty."
    if numCards is not None and not (1 <= numCards <= 20):
        return "❌ numCards must be between 1 and 20."

    params: Dict[str, Any] = {"inputText": inputText.strip()}
    if textMode               is not None: params["textMode"]               = textMode
    if numCards               is not None: params["numCards"]               = numCards
    if additionalInstructions is not None: params["additionalInstructions"] = additionalInstructions.strip()

    mode, value, error = await _start_generation(GAMMA_GENERATIONS, params, "new presentation")

    if mode == "error":
        return f"❌ Failed to start generation.\nError: {error}"
    if mode == "direct":
        return f"✅ Presentation ready!\n🔗 View here: {value}"

    # mode == "poll"
    result = await _poll_generation(value)
    return _format_result(result, "Presentation")


# ═══════════════════════════════════════════════════════════════════════════
# TOOL 2 — Generate from Template
# ═══════════════════════════════════════════════════════════════════════════
@mcp.tool()
async def generate_from_template(
    gammaId: str,
    prompt: str,
    themeId: Optional[str] = None,
    folderIds: Optional[List[str]] = None,
    exportAs: Optional[Literal["pdf", "pptx"]] = None,
    imageModel: Optional[str] = None,
    imageStyle: Optional[str] = None,
    workspaceAccess: Optional[Literal[
        "noAccess", "view", "comment", "edit", "fullAccess"
    ]] = None,
    externalAccess: Optional[Literal[
        "noAccess", "view", "comment", "edit"
    ]] = None,
    emailRecipients: Optional[List[str]] = None,
    emailAccess: Optional[Literal[
        "view", "comment", "edit", "fullAccess"
    ]] = None,
) -> str:
    """
    Generate a new presentation based on an existing Gamma template.
    The template's structure and design are preserved; content is adapted
    according to your prompt.

    Args:
        gammaId          : ID of the source Gamma template (REQUIRED).
                           Find it in the deck URL on gamma.app.
        prompt           : How to adapt the template. Can include new topic,
                           content, image URLs to embed, tone/audience guidance.
                           (REQUIRED)
        themeId          : Override the template's theme (optional).
                           Use list_themes() to find available IDs.
        folderIds        : Folder IDs to save the output into (optional).
                           Use list_folders() to find available IDs.
        exportAs         : Also export as 'pdf' or 'pptx' (optional).
                           ⚠️ Export links expire — download quickly.
        imageModel       : AI image model override (optional).
                           Examples: 'flux-1-pro', 'imagen-4-pro'
                           Check Gamma docs for the full current list.
        imageStyle       : Image style hint, max 500 chars (optional).
                           Example: 'photorealistic' or 'minimal line art'
        workspaceAccess  : Workspace member access level (optional).
        externalAccess   : Public access level (optional).
        emailRecipients  : Emails to share the result with (optional).
        emailAccess      : Access level for email recipients (optional).

    Returns:
        Link to the generated presentation and any export links.
    """
    print(f"\n🚀 generate_from_template()", flush=True)
    print(f"   gammaId  : {gammaId}", flush=True)
    print(f"   prompt   : {prompt[:100]}", flush=True)
    print(f"   exportAs : {exportAs}", flush=True)

    err = _check_api_key()
    if err:
        return err

    if not gammaId or not gammaId.strip():
        return "❌ gammaId cannot be empty."
    if not prompt or not prompt.strip():
        return "❌ prompt cannot be empty."
    if imageStyle and len(imageStyle) > 500:
        return "❌ imageStyle must be 500 characters or less."

    params: Dict[str, Any] = {
        "gammaId": gammaId.strip(),
        "prompt":  prompt.strip(),
    }

    if themeId   is not None: params["themeId"]   = themeId
    if folderIds is not None: params["folderIds"]  = folderIds
    if exportAs  is not None: params["exportAs"]   = exportAs

    image_opts: Dict[str, str] = {}
    if imageModel is not None: image_opts["model"] = imageModel
    if imageStyle is not None: image_opts["style"]  = imageStyle
    if image_opts:             params["imageOptions"] = image_opts

    sharing: Dict[str, Any] = {}
    if workspaceAccess is not None: sharing["workspaceAccess"] = workspaceAccess
    if externalAccess  is not None: sharing["externalAccess"]  = externalAccess

    email_opts: Dict[str, Any] = {}
    if emailRecipients is not None: email_opts["recipients"] = emailRecipients
    if emailAccess     is not None: email_opts["access"]     = emailAccess
    if email_opts:                  sharing["emailOptions"]   = email_opts
    if sharing:                     params["sharingOptions"]  = sharing

    endpoint = f"{GAMMA_GENERATIONS}/from-template"
    mode, value, error = await _start_generation(endpoint, params, "from-template")

    if mode == "error":
        return (
            f"❌ Failed to start from-template generation.\n"
            f"Error: {error}\n\n"
            f"💡 Double-check your gammaId — find it in the deck URL on gamma.app."
        )
    if mode == "direct":
        return f"✅ Presentation ready!\n🔗 View here: {value}"

    result = await _poll_generation(value)
    return _format_result(result, "Presentation from template")


# ═══════════════════════════════════════════════════════════════════════════
# TOOL 3 — List Themes
# NOTE: The /themes endpoint is not in the provided Gamma API docs.
# This attempts the most common REST convention. If it returns 404,
# copy themeId values directly from the Gamma app instead.
# ═══════════════════════════════════════════════════════════════════════════
@mcp.tool()
async def list_themes() -> str:
    """
    List available Gamma themes in your workspace.
    Use returned themeId values in generate_presentation() or
    generate_from_template().

    Note: If this returns a 404 error, copy themeId values directly
    from the Gamma app (Settings → Themes).

    Returns:
        List of theme names and IDs, or an error with guidance.
    """
    print(f"\n🎨 list_themes()", flush=True)

    err = _check_api_key()
    if err:
        return err

    url = f"{GAMMA_API_BASE}/themes"
    print(f"🔗 GET {url}", flush=True)

    try:
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(url, headers=_headers())

        print(f"📥 HTTP {r.status_code} | {r.text[:500]}", flush=True)

        if r.status_code == 404:
            return (
                "⚠️ The /themes endpoint is not available (404).\n"
                "💡 Copy themeId values directly from the Gamma app instead."
            )
        if not r.is_success:
            return (
                f"❌ Could not fetch themes (HTTP {r.status_code}).\n"
                f"Response: {r.text[:300]}\n\n"
                f"💡 Copy themeId values directly from the Gamma app instead."
            )

        data = r.json()
        themes = (
            data
            if isinstance(data, list)
            else data.get("themes") or data.get("data") or data.get("items") or []
        )

        if not themes:
            return (
                f"⚠️ No themes returned.\nRaw response: {data}\n\n"
                f"💡 Copy themeId values directly from the Gamma app instead."
            )

        lines = ["🎨 Available Gamma Themes", "=" * 40, ""]
        for i, theme in enumerate(themes, 1):
            name     = theme.get("name") or theme.get("title") or "Unnamed"
            theme_id = theme.get("id") or theme.get("themeId") or theme.get("_id") or "unknown"
            desc     = theme.get("description") or ""
            lines.append(f"{i:2}. {name}  (ID: {theme_id})")
            if desc:
                lines.append(f"    {desc}")
            lines.append("")

        lines += [
            "=" * 40,
            "💡 Use themeId in generate_presentation() or generate_from_template()",
        ]
        return "\n".join(lines)

    except Exception as e:
        return f"❌ Error fetching themes: {type(e).__name__}: {e}"


# ═══════════════════════════════════════════════════════════════════════════
# TOOL 4 — List Folders
# NOTE: Same caveat as list_themes — endpoint not in provided docs.
# ═══════════════════════════════════════════════════════════════════════════
@mcp.tool()
async def list_folders() -> str:
    """
    List Gamma folders in your workspace.
    Use returned folder IDs in generate_from_template(folderIds=[...]).

    Note: If this returns a 404 error, copy folder IDs directly
    from the Gamma app.

    Returns:
        List of folder names and IDs, or an error with guidance.
    """
    print(f"\n📁 list_folders()", flush=True)

    err = _check_api_key()
    if err:
        return err

    url = f"{GAMMA_API_BASE}/folders"
    print(f"🔗 GET {url}", flush=True)

    try:
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(url, headers=_headers())

        print(f"📥 HTTP {r.status_code} | {r.text[:500]}", flush=True)

        if r.status_code == 404:
            return (
                "⚠️ The /folders endpoint is not available (404).\n"
                "💡 Copy folder IDs directly from the Gamma app instead."
            )
        if not r.is_success:
            return (
                f"❌ Could not fetch folders (HTTP {r.status_code}).\n"
                f"Response: {r.text[:300]}\n\n"
                f"💡 Copy folder IDs directly from the Gamma app instead."
            )

        data = r.json()
        folders = (
            data
            if isinstance(data, list)
            else data.get("folders") or data.get("data") or data.get("items") or []
        )

        if not folders:
            return (
                f"⚠️ No folders found.\nRaw response: {data}\n\n"
                f"💡 Create folders in the Gamma app first."
            )

        lines = ["📁 Gamma Folders", "=" * 40, ""]
        for i, folder in enumerate(folders, 1):
            name      = folder.get("name") or folder.get("title") or "Unnamed"
            folder_id = folder.get("id") or folder.get("folderId") or folder.get("_id") or "unknown"
            lines.append(f"{i:2}. {name}  (ID: {folder_id})")
            lines.append("")

        lines += [
            "=" * 40,
            "💡 Use folderIds in generate_from_template(folderIds=['id1', 'id2'])",
        ]
        return "\n".join(lines)

    except Exception as e:
        return f"❌ Error fetching folders: {type(e).__name__}: {e}"


# ═══════════════════════════════════════════════════════════════════════════
# TOOL 5 — Check Generation Status
# ═══════════════════════════════════════════════════════════════════════════
@mcp.tool()
async def check_generation_status(generation_id: str) -> str:
    """
    Check the current status of a generation by its ID.
    Use this if a previous call timed out but the generation may still
    be running in Gamma's backend.

    Args:
        generation_id : The ID from a previous generation
                        e.g. 'qiZhDe7r2zG88PHSp2r0U'

    Returns:
        Current status and URL if completed.
    """
    print(f"\n🔍 check_generation_status({generation_id})", flush=True)

    err = _check_api_key()
    if err:
        return err

    if not generation_id or not generation_id.strip():
        return "❌ generation_id cannot be empty."

    url = f"{GAMMA_GENERATIONS}/{generation_id.strip()}"
    print(f"🔗 GET {url}", flush=True)

    try:
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(url, headers=_headers())

        print(f"📥 HTTP {r.status_code} | {r.text[:500]}", flush=True)

        if r.status_code == 404:
            return f"❌ Generation '{generation_id}' not found."
        if not r.is_success:
            return f"❌ HTTP {r.status_code}: {_extract_gamma_error(r)}"

        data   = r.json()
        status = data.get("status", "unknown").lower()

        if status in ("completed", "complete", "done", "success"):
            presentation_url = _extract_url(data)
            if not presentation_url:
                return (
                    f"✅ Generation completed but no URL was returned.\n"
                    f"🆔 ID: {generation_id}\n"
                    f"💡 Check https://gamma.app — your deck should be there."
                )
            pdf_url  = data.get("pdfUrl")
            pptx_url = data.get("pptxUrl")
            lines = [
                "✅ Generation COMPLETED!",
                "",
                f"🆔 ID        : {generation_id}",
                f"🔗 View here : {presentation_url}",
            ]
            if pdf_url:
                lines.append(f"📄 PDF       : {pdf_url}")
            if pptx_url:
                lines.append(f"📊 PPTX      : {pptx_url}")
            return "\n".join(lines)

        elif status in ("failed", "error", "cancelled", "canceled"):
            msg = data.get("message") or data.get("error") or status
            return (
                f"❌ Generation FAILED\n"
                f"🆔 ID     : {generation_id}\n"
                f"📊 Status : {status}\n"
                f"💬 Reason : {msg}"
            )

        else:
            return (
                f"⏳ Generation is {status.upper() or 'PENDING'}\n"
                f"🆔 ID     : {generation_id}\n"
                f"💡 Try again in a moment, or check https://gamma.app"
            )

    except Exception as e:
        return f"❌ Error checking status: {type(e).__name__}: {e}"


# ═══════════════════════════════════════════════════════════════════════════
# TOOL 6 — Test Connection
# ═══════════════════════════════════════════════════════════════════════════
@mcp.tool()
async def test_connection() -> str:
    """
    Health check for the MCP server and Gamma API connectivity.
    Run this first to verify everything is configured correctly.
    """
    print("\n🔍 test_connection()", flush=True)
    lines = ["🔍 Gamma MCP Server — Connection Test", "=" * 45, ""]

    if GAMMA_API_KEY:
        masked = f"{GAMMA_API_KEY[:6]}****{GAMMA_API_KEY[-4:]}"
        lines.append(f"✅ API Key         : {masked}")
    else:
        lines.append("❌ API Key         : NOT SET — add GAMMA_API_KEY to .env!")

    try:
        async with httpx.AsyncClient(timeout=10) as client:
            r = await client.get("https://public-api.gamma.app")
        lines.append(f"✅ Gamma API       : Reachable (HTTP {r.status_code})")
    except httpx.TimeoutException:
        lines.append("❌ Gamma API       : Connection timed out")
    except httpx.ConnectError:
        lines.append("❌ Gamma API       : Cannot connect")
    except Exception as e:
        lines.append(f"❌ Gamma API       : {e}")

    poll_total = POLL_MAX_ATTEMPTS * POLL_INTERVAL_SEC
    lines += [
        "",
        f"✅ MCP Server      : Running",
        f"🔗 Base URL        : {GAMMA_API_BASE}",
        f"⏳ Poll interval   : {POLL_INTERVAL_SEC}s",
        f"⏳ Max attempts    : {POLL_MAX_ATTEMPTS}",
        f"⏳ Max wait        : {poll_total}s ({poll_total // 60}m {poll_total % 60}s)",
        f"🔄 Error streak    : give up after {POLL_MAX_ERR_STREAK} consecutive errors",
        "",
        "🛠️  Available Tools:",
        "   1. generate_presentation    — New presentation from topic",
        "   2. generate_from_template   — Adapt an existing Gamma deck",
        "   3. list_themes              — Browse available themes",
        "   4. list_folders             — Browse workspace folders",
        "   5. check_generation_status  — Resume / check a past generation",
        "   6. test_connection          — This health check",
    ]

    output = "\n".join(lines)
    print(output, flush=True)
    return output



@mcp.tool()
async def extract_ppt_text(file_path: str) -> str:
    """
    Extract text from a PowerPoint file slide-by-slide.

    Args:
        file_path: Path to PPTX file

    Returns:
        Slide-by-slide content.
    """

    try:
        prs = Presentation(file_path)

        slides = []

        for i, slide in enumerate(prs.slides, start=1):
            text_items = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        text_items.append(txt)

            slides.append(f"Slide {i}:\n" + "\n".join(text_items))

        return "\n\n".join(slides)

    except Exception as e:
        return f"❌ Failed to extract PPT text: {e}"

# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════
def main():
    print("=" * 60, flush=True)
    print("🤖  Gamma Presentation MCP Server", flush=True)
    print("=" * 60, flush=True)

    if not GAMMA_API_KEY:
        print("❌ GAMMA_API_KEY missing from environment!", flush=True)
        print("   Create a .env file with:", flush=True)
        print("   GAMMA_API_KEY=sk-gamma-xxxxxxxxxx", flush=True)
        sys.exit(1)

    masked     = f"{GAMMA_API_KEY[:6]}****{GAMMA_API_KEY[-4:]}"
    port       = int(os.environ.get("PORT", 8000))
    poll_total = POLL_MAX_ATTEMPTS * POLL_INTERVAL_SEC

    print(f"✅ API Key     : {masked}", flush=True)
    print(f"🚀 Port        : {port}", flush=True)
    print(f"⏳ Poll config : {POLL_MAX_ATTEMPTS} × {POLL_INTERVAL_SEC}s = {poll_total}s max", flush=True)
    print(f"🛠️  Tools       : 6 registered", flush=True)
    print("=" * 60, flush=True)

    mcp.run(
        transport="streamable-http",
        host="0.0.0.0",
        port=port,
    )


if __name__ == "__main__":
    main()